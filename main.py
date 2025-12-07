import telebot
from telebot import types
from deep_translator import GoogleTranslator
from gtts import gTTS
from flask import Flask
from threading import Thread
import os
import io

# Libraries សម្រាប់អានឯកសារ
import PyPDF2
import docx
import openpyxl # សម្រាប់ Excel
from pptx import Presentation # សម្រាប់ PowerPoint

# ==========================================
# ១. ផ្នែក KEEP ALIVE (Server)
# ==========================================
app = Flask('')

@app.route('/')
def home():
    return "Bot is running with Commands & File Support!"

def run():
  app.run(host='0.0.0.0', port=8080)

def keep_alive():
    t = Thread(target=run)
    t.start()

# ==========================================
# ២. ការកំណត់ BOT
# ==========================================
API_TOKEN = os.environ.get('BOT_TOKEN', '8223217940:AAH1tHD72PojpV0f4VIkzTnUwePpyxuL9Og') 
bot = telebot.TeleBot(API_TOKEN)

BANNER_IMAGE_URL = "https://blogger.googleusercontent.com/img/b/R29vZ2xl/AVvXsEili8wBjfGex4X3AizfLATOq2G3joXpZRM15nZv191_dIZfSwhe3_0dolOpI5hmCvl9epB65IcvDOj4aeqt2cyo8PtwUEgvzOm2BIdvQhvp2QY4HT0-MgGpwPznPnfYxHuTmA8JbPlJU6TdYRuaF80qEqgFHoPISsYLRA2IQbwpdzN5WJJorx9cAkv2FQXQ/s16000/photo_2025-12-07_00-59-36.jpg"

user_preferences = {} 

LANGUAGES_MAP = {
    'km': '🇰🇭 ខ្មែរ',
    'en': '🇬🇧 អង់គ្លេស',
    'ja': '🇯🇵 ជប៉ុន',
    'ko': '🇰🇷 កូរ៉េ',
    'hi': '🇮🇳 ឥណ្ឌា',
    'zh-CN': '🇨🇳 ចិន',
    'fr': '🇫🇷 បារាំង',
    'ms': '🇲🇾 ម៉ាឡេស៊ី',
    'my': '🇲🇲 មីយ៉ាន់ម៉ា',
    'id': '🇮🇩 ឥណ្ឌូនេស៊ី',
    'ru': '🇷🇺 រុស្ស៊ី',
}

# ==========================================
# ៣. DASHBOARD & KEYBOARDS
# ==========================================
def get_main_dashboard():
    markup = types.InlineKeyboardMarkup(row_width=2)
    btn_translate = types.InlineKeyboardButton("🔤 បកប្រែអក្សរ", callback_data='menu_translate')
    btn_file = types.InlineKeyboardButton("📂 បកប្រែឯកសារ", callback_data='menu_file')
    btn_voice = types.InlineKeyboardButton("🗣️ សំឡេង", callback_data='menu_voice')
    btn_info = types.InlineKeyboardButton("ℹ️ អំពី Bot", callback_data='menu_info')
    markup.add(btn_translate, btn_file, btn_voice, btn_info)
    return markup

def get_language_keyboard():
    markup = types.InlineKeyboardMarkup(row_width=3)
    buttons = []
    for code, name in LANGUAGES_MAP.items():
        buttons.append(types.InlineKeyboardButton(name, callback_data=f'set_lang_{code}'))
    markup.add(*buttons)
    markup.add(types.InlineKeyboardButton("🔙 ត្រឡប់ទៅ Dashboard", callback_data='back_home'))
    return markup

def get_back_home_btn():
    markup = types.InlineKeyboardMarkup()
    markup.add(types.InlineKeyboardButton("🔙 ត្រឡប់ទៅ Dashboard", callback_data='back_home'))
    return markup

# ==========================================
# ៤. មុខងារជំនួយ (HELPER FUNCTIONS)
# ==========================================
def split_message(text, limit=4000):
    return [text[i:i+limit] for i in range(0, len(text), limit)]

def translate_and_reply(message, text_to_translate):
    chat_id = message.chat.id
    dest_lang = user_preferences.get(chat_id, 'km')
    
    try:
        translated = GoogleTranslator(source='auto', target=dest_lang).translate(text_to_translate)
        
        chunks = split_message(translated)
        bot.reply_to(message, f"✅ **លទ្ធផលបកប្រែ ({LANGUAGES_MAP.get(dest_lang)}):**", parse_mode='Markdown')
        
        for chunk in chunks:
            bot.send_message(chat_id, chunk)
            
        # Voice Feature
        if len(translated) < 500: 
            try:
                tts_lang = dest_lang
                if dest_lang == 'zh-CN': tts_lang = 'zh'
                tts = gTTS(text=translated, lang=tts_lang)
                filename = f"voice_{chat_id}.mp3"
                tts.save(filename)
                with open(filename, 'rb') as audio:
                    bot.send_voice(chat_id, audio)
                os.remove(filename)
            except:
                pass
            
    except Exception as e:
        bot.reply_to(message, f"⚠️ បញ្ហា៖ {e}")

# ==========================================
# ៥. មុខងារអានឯកសារ (FILE READERS)
# ==========================================
def read_file_content(file_bytes, file_ext):
    text = ""
    # 1. PDF
    if file_ext == '.pdf':
        read_pdf = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        for page in read_pdf.pages:
            text += page.extract_text() + "\n"
    # 2. Word (.docx)
    elif file_ext == '.docx':
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
    # 3. Excel (.xlsx)
    elif file_ext == '.xlsx':
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        for sheet in wb.worksheets:
            text += f"\n--- Sheet: {sheet.title} ---\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) for cell in row if cell is not None])
                text += row_text + "\n"
    # 4. PowerPoint (.pptx)
    elif file_ext == '.pptx':
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    # 5. Text File (.txt)
    elif file_ext == '.txt':
        text = file_bytes.decode('utf-8')
    return text

# ==========================================
# ៦. COMMAND HANDLERS (ដោះស្រាយបញ្ហា /help, /lang)
# ==========================================

# --- /start ---
@bot.message_handler(commands=['start'])
def send_welcome(message):
    if message.chat.id not in user_preferences:
        user_preferences[message.chat.id] = 'km'
    
    caption_text = (
        f"សួស្តី/ជម្រាបសួរបង **{message.from_user.last_name}**! 👋\n\n"
        "ខ្ញុំគឺជាមនុស្សយន្ត **សម្រាប់បកប្រែភាសារ ដែលបង្កើតដោយបង ស្រ៊ាង ស៊ីណាន** \n"
        "ខ្ញុំអាចអាន និងបកប្រែឯកសារជាច្រើនប្រភេទ៖\n"
        "📄 Word, 📕 PDF, 📊 Excel, 📽️ PPT, 📝 Text\n\n"
        "👇 **សូមជ្រើសរើសមុខងារខាងក្រោម៖**"
    )
    bot.send_photo(
        message.chat.id, 
        BANNER_IMAGE_URL, 
        caption=caption_text, 
        parse_mode='Markdown', 
        reply_markup=get_main_dashboard()
    )

# --- /help ---
@bot.message_handler(commands=['help'])
def send_help(message):
    help_text = (
        "ℹ️ **ជំនួយ (Help)**\n\n"
        "1. **Start**: វាយ /start ដើម្បីបើក Menu ដើម។\n"
        "2. **ប្តូរភាសា**: វាយ /lang ឬចុចលើប៊ូតុងក្នុង Menu។\n"
        "3. **បកប្រែ**: គ្រាន់តែផ្ញើអក្សរ ឬឯកសារ (Word, PDF, Excel...) មក ខ្ញុំនឹងបកប្រែជូន។"
    )
    bot.reply_to(message, help_text, parse_mode='Markdown')

# --- /lang ---
@bot.message_handler(commands=['lang', 'language'])
def send_language_menu(message):
    chat_id = message.chat.id
    current_lang = LANGUAGES_MAP.get(user_preferences.get(chat_id, 'km'))
    text = f"🔤 **ប្តូរភាសា**\nភាសាបច្ចុប្បន្ន៖ **{current_lang}**\n\nសូមជ្រើសរើសភាសាដែលចង់បកប្រែទៅ៖"
    bot.send_message(chat_id, text, reply_markup=get_language_keyboard(), parse_mode='Markdown')

# ==========================================
# ៧. CALLBACK & CONTENT HANDLERS
# ==========================================
@bot.callback_query_handler(func=lambda call: True)
def handle_query(call):
    chat_id = call.message.chat.id
    
    if call.data == 'back_home':
        bot.delete_message(chat_id, call.message.message_id)
        send_welcome(call.message) 
        
    elif call.data == 'menu_translate':
        # ហៅ Function /lang មកប្រើវិញ ដើម្បីកុំឱ្យសរសេរកូដស្ទួន
        send_language_menu(call.message)
        
    elif call.data == 'menu_file':
        bot.send_message(chat_id, "📂 **បកប្រែឯកសារ**\n\nសូមផ្ញើ File មកខ្ញុំ (Word, Excel, PDF, PPT, TXT)។", reply_markup=get_back_home_btn(), parse_mode='Markdown')
        
    elif call.data == 'menu_voice':
         bot.send_message(chat_id, "🎙️ **មុខងារសំឡេង**\n\nផ្ញើអក្សរមក ខ្ញុំនឹងអានជូន។", reply_markup=get_back_home_btn())

    elif call.data.startswith('set_lang_'):
        code = call.data.split('_')[2]
        user_preferences[chat_id] = code
        bot.answer_callback_query(call.id, f"ប្តូរទៅជា {LANGUAGES_MAP.get(code)}")
        bot.send_message(chat_id, f"✅ បានកំណត់ភាសា **{LANGUAGES_MAP.get(code)}** រួចរាល់!", parse_mode='Markdown')

    elif call.data == 'menu_info':
        bot.send_message(chat_id, "🤖 **Bot Info**\nSupports: PDF, DOCX, XLSX, PPTX, TXT\nVersion:2.0\nTel:087533780\nDevelop By: Sreang_Sinan", reply_markup=get_back_home_btn())

@bot.message_handler(func=lambda message: True, content_types=['text'])
def handle_text(message):
    # បើមិនមែនជា Command ទេ ទើបយកមកបកប្រែ
    if not message.text.startswith('/'):
        translate_and_reply(message, message.text)

@bot.message_handler(content_types=['document'])
def handle_docs(message):
    chat_id = message.chat.id
    try:
        file_name = message.document.file_name
        file_ext = os.path.splitext(file_name)[1].lower()
        supported_exts = ['.pdf', '.docx', '.xlsx', '.pptx', '.txt']

        if file_ext not in supported_exts:
            bot.reply_to(message, f"⚠️ ខ្ញុំមិនស្គាល់ File ប្រភេទ `{file_ext}` ទេ។", parse_mode='Markdown')
            return

        bot.send_chat_action(chat_id, 'upload_document')
        bot.reply_to(message, f"⏳ កំពុងទាញយក និងអានឯកសារ **{file_name}**...")

        file_info = bot.get_file(message.document.file_id)
        downloaded_file = bot.download_file(file_info.file_path)
        extracted_text = read_file_content(downloaded_file, file_ext)

        if len(extracted_text.strip()) == 0:
            bot.reply_to(message, "⚠️ ឯកសារនេះទទេ ឬខ្ញុំមិនអាចអានអក្សរបាន។")
            return

        bot.reply_to(message, "✅ បានអានរួចរាល់! កំពុងបកប្រែ...")
        translate_and_reply(message, extracted_text)

    except Exception as e:
        bot.reply_to(message, f"❌ មានបញ្ហាក្នុងការអានឯកសារ៖ {e}")

# ==========================================
# ៨. RUN
# ==========================================
keep_alive()
try:
    bot.infinity_polling()
except:
    pass

